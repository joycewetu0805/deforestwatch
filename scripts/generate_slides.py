"""
Génère la présentation de soutenance DeforestWatch-DRC (.pptx).

Diaporama professionnel, thème sombre cohérent avec la plateforme, pré-rempli
à partir du projet implémenté. Les résultats chiffrés proviennent du rapport
d'entraînement (mode démo) et sont à confirmer sur données réelles.

Usage : python -m scripts.generate_slides
Sortie : docs/soutenance_deforestwatch.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.utils import synthetic
from src.utils.helpers import load_json
from src.utils.logger import get_logger

log = get_logger("generate_slides")

BASE = RGBColor(0x0A, 0x0F, 0x1C)
PANEL = RGBColor(0x13, 0x1A, 0x2B)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
GREY = RGBColor(0x94, 0xA3, 0xB8)


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BASE


def _text(slide, text, left, top, width, height, size, color=WHITE, bold=False,
          align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _bullets(slide, items, left=0.9, top=1.8, width=8.2, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = "•  " + it
        run.font.size = Pt(size)
        run.font.color.rgb = WHITE
        p.space_after = Pt(8)


def _title_slide(prs, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _text(s, title, 0.9, 0.55, 8.2, 1, 30, EMERALD, bold=True)
    if subtitle:
        _text(s, subtitle, 0.9, 1.35, 8.2, 0.6, 14, GREY)
    return s


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Titre
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _text(s, "🌍 DeforestWatch-DRC", 0.9, 2.4, 8.2, 1, 40, EMERALD, bold=True, align=PP_ALIGN.CENTER)
    _text(s, "Surveillance et prédiction de la déforestation de la forêt équatoriale\n"
             "du Bassin du Congo par imagerie satellite et Machine Learning",
          0.9, 3.5, 8.2, 1.2, 18, WHITE, align=PP_ALIGN.CENTER)
    _text(s, "Joyce A. WETUNGANI · UPC / FASI · L3 Data Science · 2026",
          0.9, 5.2, 8.2, 0.5, 14, GREY, align=PP_ALIGN.CENTER)

    # 2. Problématique
    _title_slide(prs, "Problématique", "Pourquoi ce projet ?")
    _bullets(prs.slides[-1], [
        "La RDC abrite la 2e plus grande forêt tropicale humide du monde (Bassin du Congo).",
        "Déforestation accélérée : agriculture sur brûlis, exploitation, expansion des villages.",
        "Les outils globaux (Global Forest Watch) manquent d'analyse prédictive locale.",
        "Besoin : un outil d'aide à la décision adapté au contexte congolais.",
    ])

    # 3. Objectifs
    _title_slide(prs, "Objectifs")
    _bullets(prs.slides[-1], [
        "Détecter les zones déforestées par classification d'images Sentinel-2.",
        "Quantifier la perte forestière annuelle (2015–2025).",
        "Comparer trois modèles : Random Forest, XGBoost, U-Net.",
        "Prédire les zones à risque de déforestation future.",
        "Déployer une plateforme : API + dashboard d'aide à la décision.",
    ])

    # 4. Zone d'étude
    _title_slide(prs, "Zone d'étude", "Province du Mai-Ndombe (Inongo)")
    _bullets(prs.slides[-1], [
        "Forêt équatoriale du Bassin du Congo, RDC.",
        "Centre : Inongo (-1,95° S ; 18,27° E), zone ~50 × 50 km.",
        "L'un des fronts de déforestation les plus actifs du pays.",
    ])

    # 5. Données
    _title_slide(prs, "Données et sources")
    _bullets(prs.slides[-1], [
        "Sentinel-2 (ESA) — 6 bandes multispectrales, 10 m.",
        "Hansen Global Forest Change — vérité terrain (perte forestière).",
        "SRTM — altitude, pente, aspect.",
        "OpenWeatherMap — précipitations, température.",
    ])

    # 6. Architecture
    _title_slide(prs, "Architecture technique")
    _bullets(prs.slides[-1], [
        "Collecte : Google Earth Engine.",
        "Prétraitement : indices (NDVI/EVI/NDWI/NBR), masquage nuages, PySpark.",
        "Modélisation : scikit-learn, XGBoost, TensorFlow (U-Net).",
        "Déploiement : FastAPI, Streamlit, frontend React, PostgreSQL, Docker.",
    ])

    # 7. Prétraitement
    _title_slide(prs, "Prétraitement")
    _bullets(prs.slides[-1], [
        "Composites médians de saison sèche (juin–septembre).",
        "13 features par pixel : 6 bandes + 4 indices + 3 variables topographiques.",
        "Tuiles 128×128 pour le U-Net.",
        "Split spatial 70/15/15 (anti-fuite géographique).",
    ])

    # 8. Modèles
    _title_slide(prs, "Modèles de Machine Learning")
    _bullets(prs.slides[-1], [
        "Random Forest — baseline robuste et interprétable (pixel).",
        "XGBoost — gradient boosting (pixel).",
        "U-Net (CNN) — segmentation sémantique exploitant le contexte spatial.",
        "Modèle de risque — prédiction des zones à déforester à court terme.",
    ])

    # 9. Résultats — déforestation
    _title_slide(prs, "Résultats — dynamique forestière",
                 "Résultats de démonstration (à confirmer sur données réelles)")
    stats = synthetic.yearly_statistics()
    total_loss = sum(s["forest_loss_ha"] for s in stats)
    _bullets(prs.slides[-1], [
        f"Surface forestière {stats[0]['year']} : {stats[0]['total_forest_ha']:,.0f} ha.",
        f"Surface forestière {stats[-1]['year']} : {stats[-1]['total_forest_ha']:,.0f} ha.",
        f"Perte cumulée : {total_loss:,.0f} ha "
        f"({100 * total_loss / stats[0]['total_forest_ha']:.1f} %).",
    ])

    # 10. Résultats — comparaison modèles
    _title_slide(prs, "Résultats — comparaison des modèles")
    metrics = Path("data/processed/model_metrics.json")
    items = []
    if metrics.exists():
        report = load_json(metrics)
        for c in report.get("comparison", []):
            items.append(f"{c['Modèle']} — F1-macro {c['F1-macro']}, Mean IoU {c['Mean IoU']}.")
        items.append(f"Meilleur modèle : {report.get('best_model')}.")
    else:
        items = ["Lancez `make train` puis régénérez les slides pour insérer les métriques."]
    _bullets(prs.slides[-1], items)

    # 11. Carte de risque
    _title_slide(prs, "Carte de risque de déforestation")
    _bullets(prs.slides[-1], [
        "Score de risque 0–100 par pixel de forêt.",
        "Facteurs : proximité aux routes/villages, déforestation voisine, pente.",
        "Identifie les fronts forestiers prioritaires pour la conservation.",
    ])

    # 12. Plateforme / Dashboard
    _title_slide(prs, "Plateforme déployée")
    _bullets(prs.slides[-1], [
        "API FastAPI : authentification JWT + 2FA, documentation Swagger.",
        "Dashboard Streamlit : KPIs, cartes, analyse, prédictions, back-office.",
        "Frontend React : landing page et monitoring temps réel.",
        "Conteneurisé (Docker) et testé (28 tests pytest).",
    ])

    # 13. Données réelles
    _title_slide(prs, "Du prototype aux données réelles")
    _bullets(prs.slides[-1], [
        "Mode démonstration : données synthétiques réalistes, exécutable clé en main.",
        "Couche d'abstraction (DataSource) : déposer des GeoTIFF + DEMO_MODE=false.",
        "Aucune modification de code : toute l'application bascule sur le réel.",
    ])

    # 14. Limites
    _title_slide(prs, "Limites")
    _bullets(prs.slides[-1], [
        "Couverture nuageuse équatoriale persistante.",
        "Besoin de données de vérité terrain de qualité.",
        "Résultats de démonstration à valider sur images réelles.",
    ])

    # 15. Perspectives
    _title_slide(prs, "Perspectives")
    _bullets(prs.slides[-1], [
        "Imagerie radar Sentinel-1 (insensible aux nuages).",
        "Modèles spatio-temporels (ConvLSTM, transformers).",
        "Alertes en temps quasi réel pour les gestionnaires forestiers.",
        "Extension à d'autres provinces de la forêt équatoriale.",
    ])

    # 16. Merci
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _text(s, "Merci de votre attention", 0.9, 2.8, 8.2, 1, 32, EMERALD, bold=True,
          align=PP_ALIGN.CENTER)
    _text(s, "Questions ?", 0.9, 3.9, 8.2, 0.6, 18, GREY, align=PP_ALIGN.CENTER)

    return prs


def main() -> None:
    out = Path("docs/soutenance_deforestwatch.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    n_slides = len(prs.slides._sldIdLst)
    prs.save(out)
    log.info(f"Présentation générée → {out} ({n_slides} slides)")


if __name__ == "__main__":
    main()
